import fitz  # PyMuPDF
from PIL import Image, ImageTk
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from pptx.util import Inches
import os  # 파일 상단에 추가

class PDFCaptureApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF PICKER | PDF 여러 페이지 캡처 → PPT 저장기")
        self.root.minsize(800, 600)  # 최소 창 크기 설정

        # 스타일 설정
        style = ttk.Style()
        style.configure('Custom.TButton', padding=5)
        style.configure('Title.TLabel', font=('맑은 고딕', 10, 'bold'))
        
        # 메인 프레임
        main_frame = ttk.Frame(root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # 상단 컨트롤 프레임
        control_frame = ttk.Frame(main_frame)
        control_frame.pack(fill=tk.X, pady=(0, 10))

        # 파일 열기 버튼
        self.load_button = ttk.Button(control_frame, text="PDF 열기", command=self.load_pdf, style='Custom.TButton', width=12)
        self.load_button.pack(side=tk.LEFT, padx=5)

        # 파일 이름 표시 레이블
        self.file_name_label = ttk.Label(control_frame, text="파일: 없음", style='Title.TLabel', width=30)
        self.file_name_label.pack(side=tk.LEFT, padx=5)

        # 구분선
        ttk.Separator(control_frame, orient=tk.VERTICAL).pack(side=tk.LEFT, padx=10, fill=tk.Y)

        # 페이지 이동 버튼들
        self.prev_button = ttk.Button(control_frame, text="←", command=self.prev_page, state=tk.DISABLED, style='Custom.TButton', width=3)
        self.prev_button.pack(side=tk.LEFT, padx=2)

        self.page_info = ttk.Label(control_frame, text="페이지: 0/0", width=12)
        self.page_info.pack(side=tk.LEFT, padx=2)

        self.next_button = ttk.Button(control_frame, text="→", command=self.next_page, state=tk.DISABLED, style='Custom.TButton', width=3)
        self.next_button.pack(side=tk.LEFT, padx=2)

        # 구분선
        ttk.Separator(control_frame, orient=tk.VERTICAL).pack(side=tk.LEFT, padx=10, fill=tk.Y)

        # 확대/축소 버튼들
        self.zoom_out_button = ttk.Button(control_frame, text="-", command=self.zoom_out, state=tk.DISABLED, style='Custom.TButton', width=3)
        self.zoom_out_button.pack(side=tk.LEFT, padx=2)

        self.zoom_label = ttk.Label(control_frame, text="100%", width=4)
        self.zoom_label.pack(side=tk.LEFT, padx=2)

        self.zoom_in_button = ttk.Button(control_frame, text="+", command=self.zoom_in, state=tk.DISABLED, style='Custom.TButton', width=3)
        self.zoom_in_button.pack(side=tk.LEFT, padx=2)

        self.fit_width_button = ttk.Button(control_frame, text="폭맞춤", command=self.fit_width, state=tk.DISABLED, style='Custom.TButton', width=8)
        self.fit_width_button.pack(side=tk.LEFT, padx=5)

        # 구분선
        ttk.Separator(control_frame, orient=tk.VERTICAL).pack(side=tk.LEFT, padx=10, fill=tk.Y)

        # 캡처 컨트롤 버튼들
        self.reset_button = ttk.Button(control_frame, text="영역 초기화", command=self.reset_captures, state=tk.DISABLED, style='Custom.TButton', width=12)
        self.reset_button.pack(side=tk.LEFT, padx=5)

        self.save_button = ttk.Button(control_frame, text="PPT 저장", command=self.save_ppt, state=tk.DISABLED, style='Custom.TButton', width=12)
        self.save_button.pack(side=tk.LEFT, padx=5)

        # 캡처 상태 표시
        self.capture_status = ttk.Label(control_frame, text="캡처 영역: 0/2", width=12)
        self.capture_status.pack(side=tk.LEFT, padx=5)
        
        # self.capture_guide = ttk.Label(control_frame, text="(최대 2개 영역 선택 가능)", style='Title.TLabel')
        # self.capture_guide.pack(side=tk.LEFT, padx=5)

        # 캔버스 프레임
        canvas_frame = ttk.Frame(main_frame)
        canvas_frame.pack(fill=tk.BOTH, expand=True)

        # 스크롤바 추가
        h_scrollbar = ttk.Scrollbar(canvas_frame, orient=tk.HORIZONTAL)
        h_scrollbar.pack(side=tk.BOTTOM, fill=tk.X)
        v_scrollbar = ttk.Scrollbar(canvas_frame)
        v_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # 캔버스 설정
        self.canvas = tk.Canvas(canvas_frame, cursor="cross", 
                              xscrollcommand=h_scrollbar.set,
                              yscrollcommand=v_scrollbar.set)
        self.canvas.pack(fill=tk.BOTH, expand=True)

        h_scrollbar.config(command=self.canvas.xview)
        v_scrollbar.config(command=self.canvas.yview)

        # 상태바 추가
        self.status_bar = ttk.Label(root, text="PDF 파일을 열어주세요", relief=tk.SUNKEN, anchor=tk.W)
        self.status_bar.pack(side=tk.BOTTOM, fill=tk.X)

        # 기존 변수들 초기화
        self.rect_start = None
        self.rect_id = None
        self.image_on_canvas = None
        self.image_path = None
        self.display_scale = 1.0
        self.pdf_path = None
        self.doc = None
        self.current_page_index = 0
        self.capture_data = []
        self.capture_dir = "capture"
        if not os.path.exists(self.capture_dir):
            os.makedirs(self.capture_dir)

        # 메시지 관련 변수 수정 (깜빡임 관련 변수 제거)
        self.message_id = None
        self.message_text_id = None

        # 이벤트 바인딩
        self.canvas.bind("<ButtonPress-1>", self.on_mouse_down)
        self.canvas.bind("<B1-Motion>", self.on_mouse_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_mouse_up)
        
        # 방향키 바인딩 추가
        self.root.bind("<Left>", lambda e: self.prev_page())
        self.root.bind("<Right>", lambda e: self.next_page())

        # 확대/축소 관련 변수 추가
        self.zoom_scale = 1.0
        self.original_image = None

        # 캡처 관련 변수 추가
        self.capture_colors = ["red", "blue"]  # 캡처 영역 색상
        self.max_captures = 2  # 최대 캡처 개수
        self.current_capture_index = 0  # 현재 캡처 인덱스
        self.rect_ids = []  # 캡처 영역 테두리 ID 저장

    def update_status(self, message):
        self.status_bar.config(text=message)
        # 캡처 상태 업데이트
        current_captures = sum(1 for page_idx, _ in self.capture_data if page_idx == self.current_page_index)
        self.capture_status.config(text=f"캡처 영역: {current_captures}/{self.max_captures}")

    def update_page_info(self):
        if self.doc:
            self.page_info.config(text=f"페이지: {self.current_page_index + 1}/{len(self.doc)}")
        else:
            self.page_info.config(text="페이지: 0/0")

    def prev_page(self):
        if self.current_page_index > 0:
            self.current_page_index -= 1
            self.load_page_image()

    def next_page(self):
        if self.current_page_index < len(self.doc) - 1:
            self.current_page_index += 1
            self.load_page_image()

    def load_pdf(self):
        self.pdf_path = filedialog.askopenfilename(filetypes=[("PDF 파일", "*.pdf")])
        if not self.pdf_path:
            return

        pdf_name = os.path.splitext(os.path.basename(self.pdf_path))[0]
        self.current_capture_dir = os.path.join(self.capture_dir, pdf_name)
        if not os.path.exists(self.current_capture_dir):
            os.makedirs(self.current_capture_dir)

        # 파일 이름 표시 업데이트
        self.file_name_label.config(text=f"파일: {os.path.basename(self.pdf_path)}")
        
        self.doc = fitz.open(self.pdf_path)
        self.current_page_index = 0
        self.capture_data = []
        self.load_page_image()
        self.prev_button.config(state=tk.NORMAL)
        self.next_button.config(state=tk.NORMAL)
        self.save_button.config(state=tk.NORMAL)
        self.reset_button.config(state=tk.NORMAL)  # PDF 로드 시 초기화 버튼 활성화
        self.update_page_info()
        self.update_status(f"PDF 파일 로드됨: {os.path.basename(self.pdf_path)}")
        self.zoom_scale = 1.0
        self.zoom_label.config(text="100%")
        self.zoom_out_button.config(state=tk.NORMAL)
        self.zoom_in_button.config(state=tk.NORMAL)
        self.fit_width_button.config(state=tk.NORMAL)
        
        # PDF 로드 후 안내 메시지 표시 시작
        self.show_capture_message()

    def load_page_image(self):
        page = self.doc.load_page(self.current_page_index)
        pix = page.get_pixmap()
        self.image_path = os.path.join(self.current_capture_dir, "page_preview.png")
        pix.save(self.image_path)

        img = Image.open(self.image_path)
        self.original_image = img.copy()  # 원본 이미지 저장
        
        # 초기 크기를 창 크기에 맞게 조정
        canvas_width = self.canvas.winfo_width()
        canvas_height = self.canvas.winfo_height()
        
        # 이미지 비율 유지하면서 창 크기에 맞게 조정
        img_ratio = img.width / img.height
        canvas_ratio = canvas_width / canvas_height
        
        # display_scale을 매번 새로 계산
        if img_ratio > canvas_ratio:
            # 이미지가 더 넓은 경우
            self.display_scale = canvas_width / img.width
        else:
            # 이미지가 더 높은 경우
            self.display_scale = canvas_height / img.height
        
        # 이미지 표시 업데이트
        self.update_image_display()
        
        # 현재 페이지의 캡처 영역이 있는지 확인
        current_page_captures = [rect for page_idx, rect in self.capture_data if page_idx == self.current_page_index]
        if not current_page_captures:
            # 캡처 영역이 없으면 안내 메시지 표시
            self.show_capture_message()
        else:
            # 캡처 영역이 있으면 메시지 숨기기
            if self.message_id:
                self.canvas.delete(self.message_id)
                self.message_id = None
            if self.message_text_id:
                self.canvas.delete(self.message_text_id)
                self.message_text_id = None

    def update_image_display(self):
        if self.original_image:
            # 원본 이미지를 현재 확대/축소 비율로 조정
            img = self.original_image.copy()
            new_width = int(img.width * self.display_scale * self.zoom_scale)
            new_height = int(img.height * self.display_scale * self.zoom_scale)
            img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            self.tk_image_tk = ImageTk.PhotoImage(img)
            self.canvas.delete("all")
            self.image_on_canvas = self.canvas.create_image(0, 0, anchor="nw", image=self.tk_image_tk)
            
            # 현재 페이지의 캡처 영역 가져오기
            current_page_captures = [(page_idx, rect) for page_idx, rect in self.capture_data 
                                   if page_idx == self.current_page_index]
            
            # 캡처 영역을 순서대로 그리기
            self.rect_ids = []  # 이전 테두리 ID 초기화
            for idx, (page_idx, rect) in enumerate(current_page_captures):
                color = self.capture_colors[idx]  # 순서대로 색상 적용
                # 확대/축소 비율을 고려하여 좌표 조정
                x0 = rect.x0 * self.display_scale * self.zoom_scale
                y0 = rect.y0 * self.display_scale * self.zoom_scale
                x1 = rect.x1 * self.display_scale * self.zoom_scale
                y1 = rect.y1 * self.display_scale * self.zoom_scale
                rect_id = self.canvas.create_rectangle(x0, y0, x1, y1, outline=color, width=2)
                self.rect_ids.append(rect_id)
            
            self.update_page_info()
            
            # 현재 페이지에 캡처 영역이 없으면 안내 메시지 표시
            if not current_page_captures:
                self.show_capture_message(new_width, new_height)
            else:
                # 캡처 영역이 있으면 메시지 숨기기
                if self.message_id:
                    self.canvas.delete(self.message_id)
                    self.message_id = None
                if self.message_text_id:
                    self.canvas.delete(self.message_text_id)
                    self.message_text_id = None

    def zoom_in(self):
        if self.zoom_scale < 3.0:  # 최대 300%까지 확대
            self.zoom_scale *= 1.2
            self.update_zoom_label()  # 확대/축소 값 표시 업데이트
            self.update_image_display()

    def zoom_out(self):
        if self.zoom_scale > 0.3:  # 최소 30%까지 축소
            self.zoom_scale /= 1.2
            self.update_zoom_label()  # 확대/축소 값 표시 업데이트
            self.update_image_display()

    def fit_width(self):
        if self.original_image:
            # 캔버스의 현재 너비에 맞게 이미지 크기 조정
            canvas_width = self.canvas.winfo_width()
            if canvas_width > 1:  # 캔버스가 실제로 생성된 후에만 실행
                self.zoom_scale = canvas_width / (self.original_image.width * self.display_scale)
                self.update_zoom_label()  # 확대/축소 값 표시 업데이트
                self.update_image_display()

    def on_mouse_down(self, event):
        # 마우스 클릭 시 메시지와 배경 제거
        if self.message_id:
            self.canvas.delete(self.message_id)
            self.message_id = None
        if hasattr(self, 'message_text_id') and self.message_text_id:
            self.canvas.delete(self.message_text_id)
            self.message_text_id = None
        
        # 현재 페이지의 캡처 개수 확인
        current_page_captures = sum(1 for page_idx, _ in self.capture_data if page_idx == self.current_page_index)
        if current_page_captures >= self.max_captures:
            messagebox.showwarning("경고", "최대 2개의 영역만 캡처할 수 있습니다.")
            return
        
        self.rect_start = (event.x, event.y)
        if self.rect_id:
            self.canvas.delete(self.rect_id)
            self.rect_id = None

    def on_mouse_drag(self, event):
        if self.rect_start:
            x0, y0 = self.rect_start
            x1, y1 = event.x, event.y
            if self.rect_id:
                self.canvas.coords(self.rect_id, x0, y0, x1, y1)
            else:
                # 첫 번째 페이지의 캡처 개수에 따라 색상 결정
                first_page_captures = sum(1 for page_idx, _ in self.capture_data if page_idx == 0)
                color = self.capture_colors[first_page_captures]
                self.rect_id = self.canvas.create_rectangle(x0, y0, x1, y1, outline=color, width=2)

    def on_mouse_up(self, event):
        if self.rect_start:
            x0, y0 = self.rect_start
            x1, y1 = event.x, event.y
            # 확대/축소 비율을 고려하여 좌표 조정
            scale = 1 / (self.display_scale * self.zoom_scale)
            rect = fitz.Rect(min(x0, x1) * scale, min(y0, y1) * scale,
                             max(x0, x1) * scale, max(y0, y1) * scale)
            
            # 현재 페이지의 캡처 개수 확인
            current_page_captures = sum(1 for page_idx, _ in self.capture_data if page_idx == self.current_page_index)
            if current_page_captures < self.max_captures:
                # 모든 페이지에 동일한 캡처 영역 저장
                for page_idx in range(len(self.doc)):
                    self.capture_data.append((page_idx, rect))
                # 캡처 영역 테두리 ID 저장
                if self.rect_id:
                    self.rect_ids.append(self.rect_id)
                self.update_status(f"캡처 영역이 모든 페이지에 적용되었습니다.")
            self.rect_start = None
            self.rect_id = None

    def save_ppt(self):
        if not self.capture_data:
            messagebox.showwarning("경고", "선택한 영역이 없습니다.")
            return

        prs = Presentation()
        
        # 각 페이지에 대해 해당 페이지의 캡처 영역 적용
        for page_index in range(len(self.doc)):
            page = self.doc.load_page(page_index)
            
            # 현재 페이지의 캡처 영역들 가져오기
            current_page_rects = [rect for p_idx, rect in self.capture_data if p_idx == page_index]
            
            # 각 캡처 영역에 대해 별도의 슬라이드 생성
            for capture_idx, rect in enumerate(current_page_rects):
                pix = page.get_pixmap(clip=rect)
                img_path = os.path.join(self.current_capture_dir, f"temp_capture_{page_index}_{capture_idx}.png")
                pix.save(img_path)

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                img = Image.open(img_path)
                width_inch = img.width / 96
                height_inch = img.height / 96
                slide.shapes.add_picture(img_path, Inches(1), Inches(1),
                                         width=Inches(width_inch), height=Inches(height_inch))

        # PPT 파일 저장
        default_ppt_name = os.path.join(self.current_capture_dir, "presentation.pptx")
        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint 파일", "*.pptx")],
            initialfile=os.path.basename(default_ppt_name)
        )
        if save_path:
            prs.save(save_path)
            messagebox.showinfo("저장 완료", f"PPT가 '{save_path}'에 저장되었습니다.")

    def show_capture_message(self, width=None, height=None):
        if self.message_id:
            self.canvas.delete(self.message_id)
        if self.message_text_id:
            self.canvas.delete(self.message_text_id)
        
        # 캔버스 크기 가져오기
        canvas_width = width if width is not None else self.canvas.winfo_width()
        canvas_height = height if height is not None else self.canvas.winfo_height()
        
        # 반투명한 검은색 배경 추가
        self.message_id = self.canvas.create_rectangle(
            0, 0, canvas_width, canvas_height,
            fill="black",
            stipple="gray50",  # 반투명 효과
            outline=""
        )
        
        # 메시지 텍스트 추가
        self.message_text_id = self.canvas.create_text(
            canvas_width/2, canvas_height/2,
            text="희망하는 캡처 영역을 선택해주세요.",
            font=("맑은 고딕", 14, "bold"),
            fill="white",  # 흰색 텍스트
            anchor="center"
        )

    def update_zoom_label(self):
        self.zoom_label.config(text=f"{int(self.zoom_scale * 100)}%")

    def reset_captures(self):
        """모든 캡처 영역을 초기화"""
        # 모든 캡처 데이터 초기화
        self.capture_data = []
        
        # 캔버스에서 모든 캡처 영역 테두리 제거
        for rect_id in self.rect_ids:
            self.canvas.delete(rect_id)
        self.rect_ids = []
        
        # 이미지 다시 표시
        self.update_image_display()
        self.update_status("모든 캡처 영역이 초기화되었습니다.")

if __name__ == "__main__":
    root = tk.Tk()
    app = PDFCaptureApp(root)
    root.mainloop()